"""
generate_presentation.py — Generates a 12-slide Capstone Presentation deck.

Slide Layout:
  1. Title Slide
  2. Problem & Objective
  3. Data Sources
  4. ETL Architecture
  5. EDA Highlights (Part 1)
  6. EDA Highlights (Part 2)
  7. Performance Metrics (Part 1)
  8. Performance Metrics (Part 2)
  9. Dashboard Screenshots (Page 1-2)
  10. Dashboard Screenshots (Page 3-4)
  11. Key Findings & Recommendations
  12. Thank You

Author: Kaushik
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pathlib import Path

BASE_DIR = Path(__file__).parent.parent
REPORT_DIR = BASE_DIR / "reports"

NAVY = RGBColor(12, 35, 64)
GOLD = RGBColor(212, 175, 55)
GREY = RGBColor(100, 100, 100)
WHITE = RGBColor(255, 255, 255)


def add_bullet(tf, text, level=0, bold=False, size=16):
    """Helper to append a bullet paragraph to a text frame."""
    p = tf.add_paragraph()
    p.text = text
    p.level = level
    p.font.size = Pt(size)
    if bold:
        p.font.bold = True
    return p


def generate_presentation():
    """Build and save a 12-slide capstone presentation."""
    prs = Presentation()

    # ===== SLIDE 1: Title =====
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Bluestock Mutual Fund Analytics"
    slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = NAVY
    slide.shapes.title.text_frame.paragraphs[0].font.bold = True
    sub = slide.placeholders[1]
    sub.text = ("Capstone Project - Final Presentation\n"
                "End-to-End Data Engineering & Quantitative Analytics\n\n"
                "Author: Kaushik  |  July 2026")

    # ===== SLIDE 2: Problem & Objective =====
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Problem Statement & Objectives"
    slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = NAVY
    tf = slide.placeholders[1].text_frame
    tf.clear()
    tf.paragraphs[0].text = "The Problem"
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.size = Pt(20)
    add_bullet(tf, "Indian MF industry manages INR 81 Lakh Crore across 1,900+ schemes and 26 Cr folios.", 1, size=15)
    add_bullet(tf, "Data is siloed across AMFI PDFs, AMC factsheets, and brokerage platforms.", 1, size=15)
    add_bullet(tf, "No unified platform exists for cross-fund, risk-adjusted performance comparison.", 1, size=15)
    add_bullet(tf, "Our Objectives", bold=True, size=20)
    add_bullet(tf, "Build an automated ETL pipeline ingesting 10 datasets into a Star Schema.", 1, size=15)
    add_bullet(tf, "Compute 6 quantitative metrics (CAGR, Sharpe, Sortino, Alpha, Beta, MaxDD) for 40 schemes.", 1, size=15)
    add_bullet(tf, "Deploy an interactive dashboard and automated reporting engine.", 1, size=15)

    # ===== SLIDE 3: Data Sources =====
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Data Sources (10 Datasets)"
    slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = NAVY
    tf = slide.placeholders[1].text_frame
    tf.clear()
    tf.paragraphs[0].text = "580,000+ records spanning Jan 2022 - May 2026"
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.size = Pt(18)
    datasets = [
        "01_fund_master.csv - 40 schemes with metadata, expense ratios, risk grades",
        "02_nav_history.csv - ~540K daily NAV values for all 40 schemes",
        "03_aum_by_fund_house.csv - Quarterly AUM for top 10 AMCs",
        "07_scheme_performance.csv - Trailing returns, Alpha, Beta, Sharpe ratios",
        "08_investor_transactions.csv - 32K+ transactions with demographics",
        "09_portfolio_holdings.csv - Sector allocations and stock weights",
        "10_benchmark_indices.csv - Nifty 50, Nifty 100, Nifty Midcap daily closes",
    ]
    for d in datasets:
        add_bullet(tf, d, 1, size=14)

    # ===== SLIDE 4: ETL Architecture =====
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "ETL Pipeline Architecture"
    slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = NAVY
    tf = slide.placeholders[1].text_frame
    tf.clear()
    tf.paragraphs[0].text = "Extract"
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.size = Pt(20)
    add_bullet(tf, "etl_pipeline.py: Automated ingestion of 10 CSVs into data/raw/ staging area.", 1, size=15)
    add_bullet(tf, "Transform", bold=True, size=20)
    add_bullet(tf, "clean_data.py: Forward-fill (ffill) for missing NAV dates, deduplication, amount validation.", 1, size=15)
    add_bullet(tf, "Expense ratios filtered to SEBI-compliant 0.1-2.5% range.", 1, size=15)
    add_bullet(tf, "Load", bold=True, size=20)
    add_bullet(tf, "db_load.py: SQLAlchemy maps CSVs to 5-table Star Schema (fact_nav, fact_transactions, dim_fund, etc.).", 1, size=15)
    add_bullet(tf, "dim_date generated dynamically spanning 2022-2026 with year/quarter/weekday flags.", 1, size=15)

    # ===== SLIDE 5: EDA Highlights Part 1 =====
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "EDA Highlights (Part 1)"
    slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = NAVY
    tf = slide.placeholders[1].text_frame
    tf.clear()
    tf.paragraphs[0].text = "Industry AUM Growth"
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.size = Pt(20)
    add_bullet(tf, "Industry AUM grew 113% from INR 38L Cr (Q1 2022) to INR 81L Cr (Q1 2025).", 1, size=15)
    add_bullet(tf, "SBI MF, ICICI Prudential, and HDFC AMC consistently held top 3 AUM positions.", 1, size=15)
    add_bullet(tf, "SIP Inflow Trends", bold=True, size=20)
    add_bullet(tf, "Monthly SIP inflows surged from INR 12,000 Cr to INR 31,000 Cr over 3 years.", 1, size=15)
    add_bullet(tf, "SIP-to-AUM ratio stabilized at ~0.4%, indicating healthy organic capital formation.", 1, size=15)
    add_bullet(tf, "NAV Return Distributions", bold=True, size=20)
    add_bullet(tf, "Daily returns follow approximately normal distributions with slight negative skew (-0.12).", 1, size=15)
    add_bullet(tf, "8 sectoral funds showed fat-tailed kurtosis > 3.0, confirming elevated tail risk.", 1, size=15)

    # ===== SLIDE 6: EDA Highlights Part 2 =====
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "EDA Highlights (Part 2)"
    slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = NAVY
    tf = slide.placeholders[1].text_frame
    tf.clear()
    tf.paragraphs[0].text = "Transaction Demographics"
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.size = Pt(20)
    add_bullet(tf, "Maharashtra, Karnataka, Gujarat = 45% of total investment volumes.", 1, size=15)
    add_bullet(tf, "Age group 25-35 shows highest SIP adoption (avg INR 8,500/month).", 1, size=15)
    add_bullet(tf, "Tier-1 cities: 62% of value; Tier-2 cities: fastest growth at 28% YoY.", 1, size=15)
    add_bullet(tf, "Fund Category Insights", bold=True, size=20)
    add_bullet(tf, "Large Cap: 34% AUM | Mid Cap: 22% | Small Cap: 18% of total AUM.", 1, size=15)
    add_bullet(tf, "Sectoral funds: highest inflow growth BUT highest redemption volatility.", 1, size=15)
    add_bullet(tf, "Net outflows observed in sectoral funds in 3 of 12 quarters.", 1, size=15)

    # ===== SLIDE 7: Performance Metrics Part 1 =====
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Performance Metrics (Part 1)"
    slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = NAVY
    tf = slide.placeholders[1].text_frame
    tf.clear()
    tf.paragraphs[0].text = "CAGR (Compounded Annual Growth Rate)"
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.size = Pt(18)
    add_bullet(tf, "Formula: (NAV_end / NAV_start) ^ (252 / N_trading_days) - 1", 1, size=15)
    add_bullet(tf, "Uses 252 trading days (not 365) to prevent 30% understatement.", 1, size=15)
    add_bullet(tf, "Sharpe Ratio (Risk-Adjusted Return)", bold=True, size=18)
    add_bullet(tf, "Formula: (Rp - Rf) / Std(Rp) x SQRT(252), Rf = 6.5% RBI repo rate proxy.", 1, size=15)
    add_bullet(tf, "Top-tier funds achieved Sharpe > 1.2, indicating strong outperformance.", 1, size=15)
    add_bullet(tf, "Sortino Ratio (Downside-Only Risk)", bold=True, size=18)
    add_bullet(tf, "Same as Sharpe but denominator uses only negative-return-day std deviation.", 1, size=15)
    add_bullet(tf, "Prevents penalization of positive volatility (upside surprise).", 1, size=15)

    # ===== SLIDE 8: Performance Metrics Part 2 =====
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Performance Metrics (Part 2)"
    slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = NAVY
    tf = slide.placeholders[1].text_frame
    tf.clear()
    tf.paragraphs[0].text = "Alpha & Beta (OLS Regression vs NIFTY 100)"
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.size = Pt(18)
    add_bullet(tf, "Beta = market volatility correlation. Beta > 1.0 = amplified market risk.", 1, size=15)
    add_bullet(tf, "Alpha = annualized excess return. Positive Alpha = genuine managerial skill.", 1, size=15)
    add_bullet(tf, "Maximum Drawdown", bold=True, size=18)
    add_bullet(tf, "Formula: min(NAV / Running_Max - 1). Measures deepest peak-to-trough decline.", 1, size=15)
    add_bullet(tf, "Sectoral funds exceeded -25% drawdown during macro corrections.", 1, size=15)
    add_bullet(tf, "Composite Scorecard (0-100)", bold=True, size=18)
    add_bullet(tf, "Weighting: 3Yr CAGR (30%) + Sharpe (25%) + Alpha (20%) + Expense Ratio (15% inv.) + MaxDD (10% inv.)", 1, size=14)
    add_bullet(tf, "Unified percentile-based ranking enabling direct cross-fund comparison.", 1, size=15)

    # ===== SLIDE 9: Dashboard (Pages 1-2) =====
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Interactive Dashboard (Pages 1-2)"
    slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = NAVY
    tf = slide.placeholders[1].text_frame
    tf.clear()
    tf.paragraphs[0].text = "Built with Streamlit + Plotly (Alternative to Power BI)"
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.size = Pt(18)
    add_bullet(tf, "Page 1 - Industry Overview", bold=True, size=17)
    add_bullet(tf, "KPI cards: Total AUM (INR 81.2L Cr), SIP Inflows (INR 31K Cr), Folios (26.12 Cr), Schemes (1,908).", 1, size=14)
    add_bullet(tf, "Interactive line chart: Industry AUM trend 2022-2025.", 1, size=14)
    add_bullet(tf, "Bar chart: AUM by AMC with color-coded fund houses.", 1, size=14)
    add_bullet(tf, "Page 2 - Fund Performance", bold=True, size=17)
    add_bullet(tf, "Scatter plot: Return (X) vs Risk/StdDev (Y), bubble size = AUM.", 1, size=14)
    add_bullet(tf, "Sortable fund scorecard table with all metrics.", 1, size=14)
    add_bullet(tf, "Normalized NAV line (base 100) vs Nifty 50 benchmark.", 1, size=14)
    add_bullet(tf, "Dynamic slicers: Fund House, Category, Plan Type.", 1, size=14)

    # ===== SLIDE 10: Dashboard (Pages 3-4) =====
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Interactive Dashboard (Pages 3-4)"
    slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = NAVY
    tf = slide.placeholders[1].text_frame
    tf.clear()
    tf.paragraphs[0].text = "Page 3 - Investor Analytics"
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.size = Pt(18)
    add_bullet(tf, "Horizontal bar chart: Transaction amount by state.", 1, size=15)
    add_bullet(tf, "Donut chart: SIP / Lumpsum / Redemption split by amount.", 1, size=15)
    add_bullet(tf, "Bar chart: Age group vs average SIP amount.", 1, size=15)
    add_bullet(tf, "Monthly transaction volume trend line. Slicers: State, Age, City Tier.", 1, size=15)
    add_bullet(tf, "Page 4 - SIP & Market Trends", bold=True, size=18)
    add_bullet(tf, "Dual-axis chart: SIP inflow (bars) + Nifty 50 (line) over 2022-2025.", 1, size=15)
    add_bullet(tf, "Category flow heatmap: Net inflows by category across 6 months.", 1, size=15)
    add_bullet(tf, "Top 5 categories by net inflow for FY25.", 1, size=15)

    # ===== SLIDE 11: Key Findings =====
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Key Findings & Recommendations"
    slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = NAVY
    tf = slide.placeholders[1].text_frame
    tf.clear()
    tf.paragraphs[0].text = "Key Findings"
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.size = Pt(20)
    add_bullet(tf, "Top 5 funds maintained Sharpe > 1.2 and positive Alpha - genuine outperformance.", 1, size=14)
    add_bullet(tf, "Sectoral funds carry 40-60% higher tail risk (VaR/CVaR) vs diversified funds.", 1, size=14)
    add_bullet(tf, "2022 cohort investors maintain 2x higher avg SIP amounts than 2024 entrants.", 1, size=14)
    add_bullet(tf, "12% of seasoned SIP investors flagged 'at-risk' with payment gaps > 35 days.", 1, size=14)
    add_bullet(tf, "Recommendations", bold=True, size=20)
    add_bullet(tf, "Diversification Advisory: Rebalance sectoral exposure toward multi-cap strategies.", 1, size=14)
    add_bullet(tf, "SIP Retention: Target at-risk investors with personalized campaigns before mandate lapse.", 1, size=14)
    add_bullet(tf, "Geographic Expansion: Digital-first onboarding in Tier-2/3 cities.", 1, size=14)

    # ===== SLIDE 12: Thank You =====
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Thank You"
    slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = NAVY
    slide.shapes.title.text_frame.paragraphs[0].font.bold = True
    sub = slide.placeholders[1]
    sub.text = ("Bluestock Mutual Fund Analytics Capstone\n\n"
                "Author: Kaushik\n"
                "GitHub: github.com/Kaushik-verse/-mutual-fund-analysis\n\n"
                "All 5 Bonus Challenges Completed\n"
                "Questions?")

    REPORT_DIR.mkdir(parents=True, exist_ok=True)
    prs.save(REPORT_DIR / 'Bluestock_MF_Presentation.pptx')
    print("Generated 12-slide Bluestock_MF_Presentation.pptx")


if __name__ == "__main__":
    generate_presentation()
